# app.py
# Streamlit: 제안 옵션 선택 → (요청별 인라인 미리보기) → PDF/Excel/PPT 내보내기
# + 추가: (1) RFP 업로드/방향성 입력 (2) 고객/프로젝트 히스토리
# 폴더 구조:
#   DB/
#     RFP/                ← 업로드된 RFP 원본 저장
#     proposal_result/    ← 분석 산출물 및 config.json 저장
#     proposal/           ← (선택) 템플릿 등
#     clients.db          ← SQLite (자동 생성)

import os
import io
import re
import json
import sqlite3
import platform
from pathlib import Path
from datetime import datetime
from typing import Optional, Tuple, List, Dict, Any
# 맨 위 import들 사이에 추가
# 기존 (에러 발생 라인)

# 교체
from pipeline.rfp2proposal import build_flows_from_user_inputs, extract_text_from_file
import shutil  # ← 파일/폴더 삭제용

import pandas as pd
import numpy as np
import streamlit as st
# ====== PDF (ReportLab) ======
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, PageBreak, Table, TableStyle, Flowable
)
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import mm
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from inspect import signature

# ====== PPT (python-pptx) ======
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
# 맨 위 import들 사이에 추가 (안전 import)
import os, streamlit as st
from pipeline import rfp2proposal as r2p
# 필요한 심벌을 별칭으로 고정
build_flows_from_user_inputs = r2p.build_flows_from_user_inputs
extract_text_from_file       = r2p.extract_text_from_file
# ==== 간단 로그인 유틸 ====
import hashlib

def _get_users_from_secrets() -> dict:
    """
    .streamlit/secrets.toml 예시
    [auth]
    enabled = true # 로그인 활성화 여부

    [auth.users]
    admin = "sha256:xxxxxxxx"  # 비번 sha256 해시(아래 주석 참고)
    viewer = "sha256:yyyyyyyy"

    파이썬에서 SHA256 만들기:
    >>> import hashlib; hashlib.sha256("비밀번호".encode()).hexdigest()
    """
    try:
        enabled = bool(st.secrets.get("auth", {}).get("enabled", False))
        users = dict(st.secrets.get("auth", {}).get("users", {}))
        return {"enabled": enabled, "users": users}
    except Exception:
        return {"enabled": False, "users": {}}

def _hash_pw(pw: str) -> str:
    return hashlib.sha256(pw.encode()).hexdigest()

def is_authed() -> bool:
    if "auth_user" not in st.session_state:
        st.session_state.auth_user = None
    auth_conf = _get_users_from_secrets()
    if not auth_conf["enabled"]:
        return True  # 로그인 비활성화 모드면 통과
    return st.session_state.auth_user is not None

def login_form(key_prefix="auth"):
    auth_conf = _get_users_from_secrets()
    if not auth_conf["enabled"]:
        return
    st.sidebar.markdown("### 🔐 로그인")
    with st.sidebar.form(f"{key_prefix}_login_form", clear_on_submit=False):
        uid = st.text_input("아이디", key=f"{key_prefix}_uid")
        pw  = st.text_input("비밀번호", type="password", key=f"{key_prefix}_pw")
        ok  = st.form_submit_button("로그인")
    if ok:
        users = auth_conf["users"]
        if uid in users:
            saved = users[uid]
            good = (saved.split(":",1)[1] == _hash_pw(pw)) if saved.startswith("sha256:") else (saved == pw)
            if good:
                st.session_state.auth_user = uid
                st.sidebar.success(f"환영합니다, {uid}님!")
                st.rerun()
            else:
                st.sidebar.error("아이디 또는 비밀번호가 올바르지 않습니다.")
        else:
            st.sidebar.error("존재하지 않는 계정입니다.")

def logout_button():
    auth_conf = _get_users_from_secrets()
    if not auth_conf["enabled"]:
        return
    if st.session_state.get("auth_user"):
        if st.sidebar.button("로그아웃"):
            st.session_state.auth_user = None
            st.rerun()

def _set_env_from_secrets():
    # st.secrets → os.environ 주입 (없는 건 건너뜀)
    for name in ("OPENAI_API_KEY", "PERPLEXITY_API_KEY"):
        if name in st.secrets and st.secrets[name] and not os.getenv(name):
            os.environ[name] = str(st.secrets[name])

_set_env_from_secrets()

# =====================================================================================
# 전역 경로/DB 경로 (자동 생성)
# =====================================================================================
ROOT = Path(__file__).resolve().parent
DB_DIR = ROOT / "DB"
RFP_DIR = DB_DIR / "RFP"
RESULT_DIR = DB_DIR / "proposal_result"
PROPOSAL_DIR = DB_DIR / "proposal"
SQLITE_PATH = DB_DIR / "clients.db"

for p in [DB_DIR, RFP_DIR, RESULT_DIR, PROPOSAL_DIR]:
    p.mkdir(parents=True, exist_ok=True)

def _ensure_parent(path: Path):
    path.parent.mkdir(parents=True, exist_ok=True)

def _ts() -> str:
    return datetime.utcnow().strftime("%Y%m%dT%H%M%SZ")

_SAFE_CHARS = re.compile(r"[^A-Za-z0-9._-]+")
def _sanitize(name: str) -> str:
    return _SAFE_CHARS.sub("-", name.strip().replace(" ", "-"))

def _call_build_flows_adapted(fn, **kwargs):
    sig = signature(fn)
    allowed = set(sig.parameters.keys())
    # 안전 기본값: 모델 고정
    if "model_main" in allowed and "model_main" not in kwargs:
        kwargs["model_main"] = "gpt-5"
    if "model_deck" in allowed and "model_deck" not in kwargs:
        kwargs["model_deck"] = "gpt-5"
    if "out_dir" in allowed and "out_dir" not in kwargs:
        # 프로젝트 out_dir을 외부에서 넣으면 그대로 사용
        pass
    # 존재하지 않는 키는 자동 필터링 (logf 포함)
    clean = {k: v for k, v in kwargs.items() if k in allowed}
    return fn(**clean)
# =====================================================================================
# SQLite (고객/프로젝트/RFP 파일)
# =====================================================================================
def _get_conn():
    conn = sqlite3.connect(str(SQLITE_PATH))
    conn.execute("PRAGMA foreign_keys = ON;")
    return conn

def init_db():
    conn = _get_conn()
    cur = conn.cursor()
    cur.execute("""
    CREATE TABLE IF NOT EXISTS clients(
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        name TEXT UNIQUE NOT NULL,
        capability TEXT,
        headcount TEXT,
        past_projects TEXT,
        created_at TEXT NOT NULL
    );
    """)
    cur.execute("""
    CREATE TABLE IF NOT EXISTS projects(
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        client_id INTEGER NOT NULL,
        title TEXT NOT NULL,
        direction TEXT NOT NULL,
        status TEXT DEFAULT 'NEW',
        created_at TEXT NOT NULL,
        FOREIGN KEY(client_id) REFERENCES clients(id) ON DELETE CASCADE
    );
    """)
    cur.execute("""
    CREATE TABLE IF NOT EXISTS rfp_files(
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        project_id INTEGER NOT NULL,
        filename TEXT NOT NULL,
        stored_path TEXT NOT NULL,
        uploaded_at TEXT NOT NULL,
        FOREIGN KEY(project_id) REFERENCES projects(id) ON DELETE CASCADE
    );
    """)
    conn.commit()
    conn.close()

def upsert_client(name: str, capability: str="", headcount: str="", past_projects: str="") -> int:
    conn = _get_conn()
    cur = conn.cursor()
    now = datetime.utcnow().isoformat()
    cur.execute("SELECT id FROM clients WHERE name=?", (name,))
    row = cur.fetchone()
    if row:
        cid = row[0]
        if any([capability, headcount, past_projects]):
            cur.execute("""
                UPDATE clients SET
                    capability=COALESCE(NULLIF(?, ''), capability),
                    headcount=COALESCE(NULLIF(?, ''), headcount),
                    past_projects=COALESCE(NULLIF(?, ''), past_projects)
                WHERE id=?
            """, (capability, headcount, past_projects, cid))
    else:
        cur.execute("""
            INSERT INTO clients(name, capability, headcount, past_projects, created_at)
            VALUES (?, ?, ?, ?, ?)
        """, (name, capability, headcount, past_projects, now))
        cid = cur.lastrowid
    conn.commit()
    conn.close()
    return cid

def create_project(client_id: int, title: str, direction: str) -> int:
    conn = _get_conn()
    cur = conn.cursor()
    now = datetime.utcnow().isoformat()
    cur.execute("""
        INSERT INTO projects(client_id, title, direction, created_at)
        VALUES (?, ?, ?, ?)
    """, (client_id, title, direction, now))
    pid = cur.lastrowid
    conn.commit()
    conn.close()
    return pid
def render_excel_like_flow(
    df: pd.DataFrame,
    *,
    default_client: str = "",
    default_author: str = "",
    default_date: Optional[str] = None,
    pdf_body_size: int = 11,
    ppt_body_size: int = 16,
    key_prefix: str = "flow"
):
    """Tab1의 Excel 업로드 흐름을 그대로 재사용하는 공용 컴포넌트.
       - 요청별 옵션 선택 UI
       - 선택 결과(selected_df)로 PDF/Excel/PPT 내보내기
    """
    df = ensure_slim_schema(df)
    df = compute_option_big_titles(df)
    df["슬라이드번호"] = df["슬라이드번호"].astype(str)
    df["옵션번호"] = df["옵션번호"].astype(str)

    # 1) 고객 정보
    st.markdown("#### 1) 고객 정보")
    col_a, col_b, col_c = st.columns(3)
    with col_a:
        client_name = st.text_input("고객사", value=default_client, key=f"{key_prefix}_client")
    with col_b:
        author = st.text_input("작성팀", value=default_author, key=f"{key_prefix}_author")
    with col_c:
        today_str = datetime.now().strftime("%Y-%m-%d")
        date_str = st.text_input("작성일", value=(default_date or today_str), key=f"{key_prefix}_date")

    client_info = {"고객사": client_name, "작성팀": author, "작성일": date_str}

    # 2) 요청별 옵션 선택 (Tab1과 동일)
    st.markdown("#### 2) 요청별 옵션 선택 (아래에 즉시 미리보기)")
    req_ids = [x for x in df["요청 ID"].unique().tolist() if x not in ("COVER","CLOSING")]
    sel_map: Dict[str, str] = {}

    for rid in req_ids:
        sub = df[df["요청 ID"] == rid]
        req_title = S(sub["요청 제목"].iloc[0] if not sub.empty else rid)
        opts = sorted({o for o in sub["옵션번호"].unique().tolist() if S(o).isdigit()},
                      key=lambda x: int(x) if S(x).isdigit() else 999)
        if not opts:
            continue

        big_title_map = {}
        for o in opts:
            g = sub[sub["옵션번호"] == o]
            bt = S(g["옵션대제목"].iloc[0]) if not g.empty and "옵션대제목" in g.columns else ""
            big_title_map[o] = bt

        st.markdown(f"**[{rid}] {req_title}**")
        sel = st.radio(
            "옵션을 선택하세요",
            options=opts,
            horizontal=True,
            index=0,
            key=f"{key_prefix}_sel_{rid}",
            format_func=lambda o: f"{o} — {big_title_map.get(o, '')}" if big_title_map.get(o, "") else f"{o}"
        )
        sel_map[rid] = sel

        with st.container():
            render_inline_preview(rid, sub, sel)
        st.divider()

    # 3) 선택 데이터 집계
    frames = []
    cover_rows = df[df["요청 ID"] == "COVER"]
    if not cover_rows.empty:
        frames.append(cover_rows)
    for rid, sel in sel_map.items():
        sub = df[df["요청 ID"] == rid]
        overview = sub[sub["슬라이드번호"] == "OVERVIEW"]
        if not overview.empty:
            frames.append(overview)
        part = sub[sub["옵션번호"] == sel]
        frames.append(part)
    closing_rows = df[df["요청 ID"] == "CLOSING"]
    if not closing_rows.empty:
        frames.append(closing_rows)
    selected_df = pd.concat(frames, ignore_index=True) if frames else pd.DataFrame(columns=df.columns)

    # 4) 내보내기 (Tab1과 동일)
    st.markdown("#### 3) 내보내기")
    c1, c2, c3 = st.columns(3)
    with c1:
        if st.button("📄 PDF 생성", use_container_width=True, key=f"{key_prefix}_pdf"):
            try:
                pdf_bytes = build_pdf(selected_df, client_info, body_size=pdf_body_size)
                st.success("PDF 생성 완료")
                st.download_button("PDF 다운로드", data=pdf_bytes, file_name="proposal_options_auto.pdf",
                                   mime="application/pdf", use_container_width=True)
            except Exception as e:
                st.error(f"PDF 생성 오류: {e}")
    with c2:
        if st.button("📊 Excel 생성", use_container_width=True, key=f"{key_prefix}_xlsx"):
            try:
                xlsx_bytes = build_excel(selected_df)
                st.success("Excel 생성 완료")
                st.download_button("Excel 다운로드", data=xlsx_bytes, file_name="proposal_options_auto.xlsx",
                                   mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                                   use_container_width=True)
            except Exception as e:
                st.error(f"Excel 생성 오류: {e}")
    with c3:
        if st.button("🖼️ PPT 생성", use_container_width=True, key=f"{key_prefix}_ppt"):
            try:
                ppt_bytes = build_ppt(selected_df, client_info, body_size=ppt_body_size)
                st.success("PPT 생성 완료")
                st.download_button("PPT 다운로드", data=ppt_bytes, file_name="proposal_options_auto.pptx",
                                   mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                   use_container_width=True)
            except Exception as e:
                st.error(f"PPT 생성 오류: {e}")

    with st.expander("선택 데이터 미리보기(전체)", expanded=False):
        st.dataframe(selected_df, use_container_width=True)

def attach_rfp(project_id: int, filename: str, stored_path: Path):
    conn = _get_conn()
    cur = conn.cursor()
    now = datetime.utcnow().isoformat()
    cur.execute("""
        INSERT INTO rfp_files(project_id, filename, stored_path, uploaded_at)
        VALUES (?, ?, ?, ?)
    """, (project_id, filename, str(stored_path), now))
    conn.commit()
    conn.close()

def fetch_client_names() -> List[str]:
    conn = _get_conn()
    rows = conn.execute("SELECT name FROM clients ORDER BY name").fetchall()
    conn.close()
    return [r[0] for r in rows]

def fetch_client_id_by_name(name: str) -> Optional[int]:
    conn = _get_conn()
    row = conn.execute("SELECT id FROM clients WHERE name=?", (name,)).fetchone()
    conn.close()
    return row[0] if row else None

def fetch_client_info(client_id: int) -> Tuple[str, str, str]:
    conn = _get_conn()
    row = conn.execute("SELECT capability, headcount, past_projects FROM clients WHERE id=?", (client_id,)).fetchone()
    conn.close()
    return row if row else ("", "", "")

def list_projects(client_id: int):
    conn = _get_conn()
    rows = conn.execute("""
        SELECT id, title, direction, created_at, status
        FROM projects WHERE client_id=? ORDER BY id DESC
    """, (client_id,)).fetchall()
    conn.close()
    return rows

def list_rfp_files(project_id: int):
    conn = _get_conn()
    rows = conn.execute("""
        SELECT filename, stored_path, uploaded_at
        FROM rfp_files WHERE project_id=? ORDER BY id DESC
    """, (project_id,)).fetchall()
    conn.close()
    return rows
def _proj_result_paths(project_id: int) -> Dict[str, Path]:
    base = RESULT_DIR / str(project_id)
    return {
        "config": base / "config.json",
        "slim_json": base / "slim_master_slide_flows.json",
        "slim_xlsx": base / "slim_master_slide_flows.xlsx",
        "auto_df_xlsx": base / f"auto_df_{project_id}.xlsx",
    }

def _safe_read_json(path: Path) -> dict:
    try:
        if path.exists():
            return json.loads(path.read_text(encoding="utf-8"))
    except Exception:
        pass
    return {}

def _show_file_download(path: Path, label: str, key=None):
    if path.exists():
        st.download_button(
            label=f"📥 {label}",
            data=path.read_bytes(),
            file_name=path.name,
            key=key or f"dl_{label}_{path.name}"
        )
    else:
        st.caption(f"⛔ {label} 파일 없음")

def _show_project_full_result(project_id: int, key_prefix: str = "pview"):
    paths = _proj_result_paths(project_id)

    st.markdown("#### 결과물 다운로드")
    c1, c2, c3, c4 = st.columns(4)
    with c1: _show_file_download(paths["slim_xlsx"], "슬림 플로우 Excel", key=f"{key_prefix}_slim_xlsx")
    with c2: _show_file_download(paths["slim_json"], "슬림 플로우 JSON", key=f"{key_prefix}_slim_json")
    with c3: _show_file_download(paths["auto_df_xlsx"], "자동 DF Excel", key=f"{key_prefix}_auto_xlsx")
    with c4: _show_file_download(paths["config"], "config.json", key=f"{key_prefix}_config")
    st.markdown("---")

    st.markdown("#### 업로드된 RFP")
    rfp_files = list_rfp_files(project_id)
    if not rfp_files:
        st.write("- (없음)")
    else:
        for i, (fn, sp, up) in enumerate(rfp_files, start=1):
            st.code(f"{fn} | {sp} | {up}", language="text")
            p = Path(sp)
            if p.exists():
                st.download_button("RFP 다운로드", data=p.read_bytes(), file_name=Path(fn).name, key=f"{key_prefix}_rfp_{i}")

    st.markdown("---")
    st.markdown("#### 자동 DF 미리보기")
    if paths["auto_df_xlsx"].exists():
        try:
            df = pd.read_excel(paths["auto_df_xlsx"])
            st.dataframe(df.head(200), use_container_width=True, key=f"{key_prefix}_auto_df_table")
        except Exception as e:
            st.caption(f"자동 DF 읽기 실패: {e}")
    else:
        st.caption("자동 DF 파일이 없습니다.")

    st.markdown("#### 슬림 플로우 Excel 미리보기")
    if paths["slim_xlsx"].exists():
        try:
            df2 = pd.read_excel(paths["slim_xlsx"])
            st.dataframe(df2.head(300), use_container_width=True, key=f"{key_prefix}_slim_df_table")
        except Exception as e:
            st.caption(f"슬림 플로우 Excel 읽기 실패: {e}")
    else:
        st.caption("슬림 플로우 Excel 파일이 없습니다.")

    st.markdown("#### 슬림 플로우 JSON 요약")
    m = _safe_read_json(paths["slim_json"])
    if m:
        sec_cnt = len(m.get("sections", []))
        st.write(f"- 섹션 수: **{sec_cnt}**")
        if sec_cnt:
            ex = m["sections"][0]
            st.write("- 예시 섹션 제목:", ex.get("req_title") or ex.get("title") or "(제목 없음)")
            st.json({"cover": m.get("cover", {}), "first_section_overview": ex.get("overview_slide", {})})
    else:
        st.caption("요약 불가 (파일 없음 또는 파싱 실패)")

# ================================
# 삭제 헬퍼 (프로젝트/클라이언트)
# ================================
def delete_project(project_id: int) -> None:
    """프로젝트와 연결된 RFP/결과물 폴더를 안전 삭제 후, DB 레코드 제거."""
    # 파일/폴더 정리
    try:
        # 업로드된 RFP 폴더 제거
        proj_rfp_dir = RFP_DIR / str(project_id)
        shutil.rmtree(proj_rfp_dir, ignore_errors=True)
        # 결과물 폴더 제거
        proj_out_dir = RESULT_DIR / str(project_id)
        shutil.rmtree(proj_out_dir, ignore_errors=True)
    except Exception:
        pass

    # DB 삭제 (rfp_files는 FK ON DELETE CASCADE로 함께 삭제)
    conn = _get_conn()
    try:
        conn.execute("DELETE FROM projects WHERE id=?", (project_id,))
        conn.commit()
    finally:
        conn.close()


def delete_client_and_all(client_id: int) -> None:
    """클라이언트 전체 삭제 (모든 프로젝트 + 파일 포함)."""
    # 먼저 해당 클라이언트의 프로젝트 폴더들 삭제
    conn = _get_conn()
    try:
        pids = [r[0] for r in conn.execute("SELECT id FROM projects WHERE client_id=?", (client_id,)).fetchall()]
    finally:
        conn.close()
    for pid in pids:
        delete_project(pid)

    # 클라이언트 레코드 삭제
    conn = _get_conn()
    try:
        conn.execute("DELETE FROM clients WHERE id=?", (client_id,))
        conn.commit()
    finally:
        conn.close()

# =====================================================================================
# 안전 문자열/파싱 유틸 (기존)
# =====================================================================================
def S(x: Any) -> str:
    if x is None:
        return ""
    if isinstance(x, float) and np.isnan(x):
        return ""
    return str(x)
# 🔧 NEW: DF 스키마 보강 함수
REQUIRED_SLIM_COLS = ["요청 ID","요청 제목","옵션번호","슬라이드번호","제목","부제목","본문초안",
    "왜_이_옵션","적합_시그널","리스크","완화책","타임라인","URL","옵션대제목"]

def ensure_slim_schema(df: pd.DataFrame) -> pd.DataFrame:
    df=df.copy()
    for c in REQUIRED_SLIM_COLS:
        if c not in df.columns: df[c]=""
    df["슬라이드번호"]=df["슬라이드번호"].astype(str)
    df["옵션번호"]=df["옵션번호"].astype(str)
    return df
def parse_url_list(val: Any) -> List[str]:
    if val is None or (isinstance(val, float) and np.isnan(val)):
        return []
    if isinstance(val, list):
        return [S(u).strip() for u in val if S(u).strip()]
    s = S(val)
    parts = []
    for tok in s.replace(";", "\n").split("\n"):
        for sub in tok.split(","):
            u = sub.strip()
            if u:
                parts.append(u)
    return parts

def parse_timeline(val: Any) -> List[Dict[str, Any]]:
    if val is None or (isinstance(val, float) and np.isnan(val)):
        return []
    if isinstance(val, list):
        return val
    s = S(val).strip()
    if not s:
        return []
    try:
        obj = json.loads(s)
        if isinstance(obj, list):
            return obj
    except Exception:
        pass
    return []

def try_extract_overview_table_from_row(row: pd.Series) -> Optional[Dict[str, Any]]:
    keys = ("columns", "rows")
    for col in row.index:
        v = row[col]
        if isinstance(v, dict) and all(k in v for k in keys):
            return v
        s = S(v).strip()
        if s.startswith("{") and s.endswith("}"):
            try:
                obj = json.loads(s)
                if isinstance(obj, dict) and all(k in obj for k in keys):
                    return obj
            except Exception:
                continue
    return None

# =====================================================================================
# 텍스트 정리 유틸 (기존)
# =====================================================================================
def strip_wrapper_quotes(s: str) -> str:
    t = s.strip()
    while (t.startswith('"') and t.endswith('"')) or (t.startswith("'") and t.endswith("'")) or (t.startswith("`") and t.endswith("`")):
        t = t[1:-1].strip()
    return t

def sanitize_text(raw: Any) -> str:
    s = S(raw)
    if not s:
        return ""
    s = s.replace("\r\n", "\n").replace("\r", "\n")
    s = strip_wrapper_quotes(s)
    cleaned_lines: List[str] = []
    for line in s.split("\n"):
        t = line.strip()
        if not t:
            continue
        if t.startswith("[붙여넣기]"):
            continue
        cleaned_lines.append(t)
    return "\n".join(cleaned_lines)

def lines_for_display(raw: Any) -> List[str]:
    s = sanitize_text(raw)
    if not s:
        return []
    return [ln.strip() for ln in s.split("\n") if ln.strip()]

def first_meaningful_line(raw: Any) -> str:
    for ln in lines_for_display(raw):
        t = ln.lstrip("-•·").strip()
        if t:
            return t
    return ""

# =====================================================================================
# 한글 폰트 자동 탐지/등록 (기존)
# =====================================================================================
def _candidate_font_paths() -> list[Tuple[str, Optional[int], str]]:
    sys = platform.system()
    cands: list[Tuple[str, Optional[int], str]] = []
    env_path = os.getenv("KOREAN_TTF_PATH")
    if env_path and os.path.exists(env_path):
        idx = None
        if env_path.lower().endswith(".ttc"):
            try:
                idx = int(os.getenv("KOREAN_TTC_INDEX", "0"))
            except:
                idx = 0
        cands.append((env_path, idx, "KR-Body"))

    if sys == "Windows":
        win_fonts = r"C:\Windows\Fonts"
        cands += [
            (os.path.join(win_fonts, "malgun.ttf"), None, "MalgunGothic-Regular"),
            (os.path.join(win_fonts, "malgunbd.ttf"), None, "MalgunGothic-Bold"),
            (os.path.join(win_fonts, "NanumGothic.ttf"), None, "NanumGothic"),
            (os.path.join(win_fonts, "NotoSansKR-Regular.otf"), None, "NotoSansKR-Regular"),
        ]
    elif sys == "Darwin":
        cands += [
            ("/Library/Fonts/AppleSDGothicNeo.ttc", 0, "AppleSDGothicNeo-0"),
            ("/System/Library/Fonts/AppleSDGothicNeo.ttc", 0, "AppleSDGothicNeo-0"),
            ("/Library/Fonts/NanumGothic.ttf", None, "NanumGothic"),
            ("/Library/Fonts/NotoSansKR-Regular.otf", None, "NotoSansKR-Regular"),
        ]
    else:
        cands += [
            ("/usr/share/fonts/truetype/nanum/NanumGothic.ttf", None, "NanumGothic"),
            ("/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc", 0, "NotoSansCJK-Regular"),
            ("/usr/share/fonts/opentype/noto/NotoSansKR-Regular.otf", None, "NotoSansKR-Regular"),
            ("/usr/share/fonts/truetype/noto/NotoSansKR-Regular.otf", None, "NotoSansKR-Regular"),
        ]
    out, seen = [], set()
    for p, idx, name in cands:
        if p and os.path.exists(p):
            key = (p, idx, name)
            if key not in seen:
                seen.add(key)
                out.append(key)
    return out

def register_korean_font_for_pdf() -> Optional[str]:
    for path, idx, name in _candidate_font_paths():
        try:
            if path.lower().endswith(".ttc"):
                TT = TTFont(name, path, subfontIndex=(0 if idx is None else idx))
            else:
                TT = TTFont(name, path)
            pdfmetrics.registerFont(TT)
            return name
        except Exception as e:
            print(f"[PDF] Font register failed: {path} -> {e}")
            continue
    return None

def _ppt_ko_font_name() -> str:
    sys = platform.system()
    return {
        "Windows": "Malgun Gothic",
        "Darwin": "Apple SD Gothic Neo",
        "Linux": "Noto Sans CJK KR"
    }.get(sys, "Malgun Gothic")

PP_KO_FONT = _ppt_ko_font_name()
PDF_KO_FONT = register_korean_font_for_pdf()

# =====================================================================================
# PDF 스타일 + HR (기존)
# =====================================================================================
def build_pdf_styles() -> Dict[str, ParagraphStyle]:
    styles = getSampleStyleSheet()
    base = PDF_KO_FONT or styles["Normal"].fontName
    if "K-Body" not in styles:
        styles.add(ParagraphStyle(
            name="K-Body", parent=styles["Normal"],
            fontName=base, fontSize=11, leading=14,
            spaceBefore=3, spaceAfter=4, textColor=colors.black
        ))
    if "K-H3" not in styles:
        styles.add(ParagraphStyle(
            name="K-H3", parent=styles["Normal"],
            fontName=base, fontSize=13, leading=16,
            spaceBefore=6, spaceAfter=3, textColor=colors.HexColor("#222")
        ))
    if "K-H2" not in styles:
        styles.add(ParagraphStyle(
            name="K-H2", parent=styles["Normal"],
            fontName=base, fontSize=15, leading=18,
            spaceBefore=8, spaceAfter=4, textColor=colors.HexColor("#111")
        ))
    if "K-H1" not in styles:
        styles.add(ParagraphStyle(
            name="K-H1", parent=styles["Normal"],
            fontName=base, fontSize=18, leading=22,
            spaceBefore=10, spaceAfter=6, textColor=colors.HexColor("#0D0D0D")
        ))
    if "K-Title" not in styles:
        styles.add(ParagraphStyle(
            name="K-Title", parent=styles["Title"],
            fontName=base, fontSize=22, leading=26,
            spaceBefore=8, spaceAfter=8, alignment=1, textColor=colors.HexColor("#0A0A0A")
        ))
    if "K-Label" not in styles:
        styles.add(ParagraphStyle(
            name="K-Label", parent=styles["Normal"],
            fontName=base, fontSize=9, leading=11, textColor=colors.HexColor("#666")
        ))
    return styles

class HR(Flowable):
    def __init__(self, width=1, thickness=0.5, color=colors.HexColor("#DDDDDD"), spaceBefore=6, spaceAfter=6):
        Flowable.__init__(self)
        self.width = width
        self.thickness = thickness
        self.color = color
        self.spaceBefore = spaceBefore
        self.spaceAfter = spaceAfter
    def wrap(self, availWidth, availHeight):
        self._w = availWidth if self.width == 1 else min(self.width, availWidth)
        return self._w, self.thickness + self.spaceBefore + self.spaceAfter
    def draw(self):
        self.canv.saveState()
        self.canv.setStrokeColor(self.color)
        self.canv.setLineWidth(self.thickness)
        self.canv.line(0, 0, self._w, 0)
        self.canv.restoreState()

PDF_STYLES = build_pdf_styles()

# =====================================================================================
# PPT 도우미 (기존)
# =====================================================================================
def apply_ppt_text_style(shape, size_pt: int = 16, bold: bool = False, align: str = "left", line_spacing: float = 1.2):
    if not hasattr(shape, "text_frame") or shape.text_frame is None:
        return
    tf = shape.text_frame
    if tf.paragraphs:
        if align == "center":
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        elif align == "right":
            tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
        else:
            tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    for p in tf.paragraphs:
        try:
            p.line_spacing = line_spacing
        except:
            pass
        for r in p.runs:
            r.font.name = PP_KO_FONT
            r.font.size = Pt(size_pt)
            r.font.bold = bool(bold)

def add_textbox(slide, left_in, top_in, width_in, height_in, text="", size=16, bold=False, align="left", line_spacing=1.2):
    tx = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = S(text)
    apply_ppt_text_style(tx, size_pt=size, bold=bold, align=align, line_spacing=line_spacing)
    return tx

def add_title_subtitle(slide, title, subtitle):
    title_box = add_textbox(slide, 0.8, 0.7, 11.0, 1.2, S(title), size=34, bold=True, align="left", line_spacing=1.1)
    if S(subtitle):
        subtitle_box = add_textbox(slide, 0.8, 1.7, 11.0, 0.8, S(subtitle), size=18, bold=False, align="left")
    else:
        subtitle_box = None
    return title_box, subtitle_box

def bullets_from_paragraphs(slide, left, top, width, height, lines: List[str], size=16):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.clear()
    first = True
    for line in lines:
        if not S(line).strip():
            continue
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = S(line)
        p.level = 0
    apply_ppt_text_style(tb, size_pt=size, bold=False, align="left", line_spacing=1.25)
    return tb

# =====================================================================================
# 옵션 대제목 생성 (기존)
# =====================================================================================
def compute_option_big_titles(df: pd.DataFrame) -> pd.DataFrame:
    df = df.copy()
    if "옵션대제목" not in df.columns:
        df["옵션대제목"] = ""
    for (rid, opt), g in df.groupby(["요청 ID", "옵션번호"]):
        if not S(opt).isdigit():
            continue
        existing = S(g["옵션대제목"].dropna().astype(str).head(1).tolist()[0] if not g["옵션대제목"].empty else "")
        if existing:
            big = existing
        else:
            big = ""
            meta = g[g["슬라이드번호"] == "META"]
            if not meta.empty:
                mt = S(meta.iloc[0].get("제목"))
                if mt:
                    big = mt
            if not big:
                detail = g[g["슬라이드번호"].apply(lambda v: S(v).isdigit())].copy()
                if not detail.empty:
                    detail["슬라이드번호"] = detail["슬라이드번호"].astype(int)
                    detail = detail.sort_values("슬라이드번호")
                    big = S(detail.iloc[0].get("제목"))
            if not big:
                big = f"옵션 {opt}"
        df.loc[(df["요청 ID"] == rid) & (df["옵션번호"] == opt), "옵션대제목"] = big
    return df

# =====================================================================================
# PDF/PPT/Excel 생성 (기존)
# =====================================================================================
def build_pdf(selected_df: pd.DataFrame, client_info: Dict[str, str], body_size=11) -> bytes:
    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf, pagesize=A4,
        leftMargin=18*mm, rightMargin=18*mm, topMargin=14*mm, bottomMargin=14*mm
    )
    styles = PDF_STYLES
    styles["K-Body"].fontSize = body_size
    styles["K-Body"].leading = int(body_size * 1.3)

    story: List[Any] = []
    title = f"{S(client_info.get('고객사',''))} 제안 옵션 패키지"
    sub = f"{S(client_info.get('작성팀',''))} · {S(client_info.get('작성일',''))}"
    story += [Spacer(1, 18), Paragraph(title, styles["K-Title"]), Paragraph(sub, styles["K-Label"]), HR()]

    summary_rows = []
    for req_id, grp in selected_df.groupby("요청 ID"):
        if req_id in ("COVER", "CLOSING"):
            continue
        req_title = S(grp["요청 제목"].iloc[0] if "요청 제목" in grp.columns else req_id)
        opt = ""
        big = ""
        opts = [x for x in grp["옵션번호"].unique().tolist() if S(x).isdigit()]
        if opts:
            opt = S(opts[0])
            g2 = grp[grp["옵션번호"] == opt]
            if not g2.empty and "옵션대제목" in g2.columns:
                big = S(g2["옵션대제목"].iloc[0])
        summary_rows.append([S(req_id), req_title, opt, big])
    if summary_rows:
        data = [["요청 ID", "요청 제목", "선택 옵션", "옵션 대제목"]] + summary_rows
        colw = [22*mm, 98*mm, 18*mm, 44*mm]
        tbl = Table(data, hAlign='LEFT', colWidths=colw)
        tbl.setStyle(TableStyle([
            ('FONTNAME', (0,0), (-1,-1), PDF_KO_FONT or 'Helvetica'),
            ('FONTSIZE', (0,0), (-1,0), 11),
            ('FONTSIZE', (0,1), (-1,-1), 9.8),
            ('BACKGROUND', (0,0), (-1,0), colors.HexColor("#F6F6F6")),
            ('ALIGN', (0,0), (-1,0), 'CENTER'),
            ('ALIGN', (0,1), (0,-1), 'CENTER'),
            ('ALIGN', (2,1), (2,-1), 'CENTER'),
            ('INNERGRID', (0,0), (-1,-1), 0.25, colors.HexColor("#DDD")),
            ('BOX', (0,0), (-1,-1), 0.5, colors.HexColor("#CCC")),
            ('TOPPADDING', (0,0), (-1,-1), 4),
            ('BOTTOMPADDING', (0,0), (-1,-1), 4),
        ]))
        story += [tbl, HR()]

    ordered_ids = [x for x in selected_df["요청 ID"].unique() if x not in ("COVER","CLOSING")]

    for idx_req, req_id in enumerate(ordered_ids):
        grp = selected_df[selected_df["요청 ID"] == req_id]
        req_title = S(grp["요청 제목"].iloc[0] if "요청 제목" in grp.columns else req_id)
        story += [Paragraph(f"[{S(req_id)}] {req_title}", styles["K-H1"]), HR()]

        sel_opts = [x for x in grp["옵션번호"].unique().tolist() if S(x).isdigit()]
        sel = S(sel_opts[0]) if sel_opts else ""
        big = ""
        if sel:
            gg = grp[grp["옵션번호"] == sel]
            if not gg.empty and "옵션대제목" in gg.columns:
                big = S(gg["옵션대제목"].iloc[0])
        if sel:
            story.append(Paragraph(f"옵션 {sel} · {big}", styles["K-H2"]))
            story.append(HR(color=colors.HexColor("#EEEEEE")))

        over = grp[grp["슬라이드번호"] == "OVERVIEW"]
        if not over.empty:
            ov = over.iloc[0]
            ov_tab = try_extract_overview_table_from_row(ov)
            if ov_tab and ov_tab.get("columns") and ov_tab.get("rows"):
                t = Table([ov_tab["columns"]] + ov_tab["rows"], hAlign='LEFT')
                t.setStyle(TableStyle([
                    ('FONTNAME', (0,0), (-1,-1), PDF_KO_FONT or 'Helvetica'),
                    ('FONTSIZE', (0,0), (-1,0), 10.5),
                    ('FONTSIZE', (0,1), (-1,-1), 9.5),
                    ('BACKGROUND', (0,0), (-1,0), colors.HexColor("#F2F2F2")),
                    ('INNERGRID', (0,0), (-1,-1), 0.25, colors.HexColor("#E1E1E1")),
                    ('BOX', (0,0), (-1,-1), 0.5, colors.HexColor("#D0D0D0")),
                    ('TOPPADDING', (0,0), (-1,-1), 3),
                    ('BOTTOMPADDING', (0,0), (-1,-1), 3),
                ]))
                story += [t, HR()]

        meta = grp[grp["슬라이드번호"] == "META"]
        if not meta.empty:
            m = meta.iloc[0]
            parts = [
                ("왜 이 옵션인가", m.get("왜_이_옵션")),
                ("적합 시그널", m.get("적합_시그널")),
                ("리스크", m.get("리스크")),
                ("완화책", m.get("완화책")),
            ]
            for (h, b) in parts:
                body_lines = lines_for_display(b)
                if not body_lines:
                    continue
                story.append(Paragraph(S(h), styles["K-H3"]))
                for ln in body_lines:
                    story.append(Paragraph(S(ln), styles["K-Body"]))
                story.append(HR())
            tl = parse_timeline(m.get("타임라인"))
            if tl:
                story += [Paragraph("타임라인(주)", styles["K-H3"])]
                tidata = [["Phase", "기간(주)"]] + [[S(x.get("phase")), S(x.get("duration_weeks"))] for x in tl]
                tt = Table(tidata, hAlign='LEFT', colWidths=[110*mm, 30*mm])
                tt.setStyle(TableStyle([
                    ('FONTNAME', (0,0), (-1,-1), PDF_KO_FONT or 'Helvetica'),
                    ('FONTSIZE', (0,0), (-1,0), 10.5),
                    ('FONTSIZE', (0,1), (-1,-1), 9.5),
                    ('BACKGROUND', (0,0), (-1,0), colors.HexColor("#F7F7F7")),
                    ('INNERGRID', (0,0), (-1,-1), 0.25, colors.HexColor("#DDDDDD")),
                    ('BOX', (0,0), (-1,-1), 0.5, colors.HexColor("#CCCCCC")),
                ]))
                story += [tt, HR(color=colors.HexColor("#EEEEEE"))]

        detail = grp[grp["슬라이드번호"].apply(lambda v: S(v).isdigit())].copy()
        if not detail.empty:
            detail["슬라이드번호"] = detail["슬라이드번호"].astype(int)
            detail = detail.sort_values(by=["슬라이드번호"])
            for j, (_, r) in enumerate(detail.iterrows()):
                story.append(Paragraph(S(r.get("제목")), styles["K-H2"]))
                if S(r.get("부제목")):
                    story.append(Paragraph(S(r.get("부제목")), styles["K-H3"]))
                for ln in lines_for_display(r.get("본문초안")):
                    story.append(Paragraph(S(ln), styles["K-Body"]))
                urls = parse_url_list(r.get("URL"))
                if urls:
                    story.append(Paragraph("참고 URL", styles["K-H3"]))
                    for u in urls:
                        story.append(Paragraph(S(u), styles["K-Body"]))
                if j < len(detail) - 1:
                    story.append(HR())

        if idx_req < len(ordered_ids) - 1:
            story.append(PageBreak())

    doc.build(story)
    return buf.getvalue()

def build_ppt(selected_df: pd.DataFrame, client_info: Dict[str, str], body_size=16) -> bytes:
    prs = Presentation()
    blank = prs.slide_layouts[6]

    cover = prs.slides.add_slide(blank)
    title = f"{S(client_info.get('고객사',''))} 제안 옵션 패키지"
    subtitle = f"{S(client_info.get('작성팀',''))} · {S(client_info.get('작성일',''))}"
    add_title_subtitle(cover, title, subtitle)

    for req_id, grp in selected_df.groupby("요청 ID"):
        if req_id in ("COVER", "CLOSING"):
            continue
        req_title = S(grp["요청 제목"].iloc[0] if "요청 제목" in grp.columns else req_id)

        sel_opts = [x for x in grp["옵션번호"].unique().tolist() if S(x).isdigit()]
        sel = S(sel_opts[0]) if sel_opts else ""
        big = ""
        if sel:
            gg = grp[grp["옵션번호"] == sel]
            if not gg.empty and "옵션대제목" in gg.columns:
                big = S(gg["옵션대제목"].iloc[0])

        s = prs.slides.add_slide(blank)
        add_title_subtitle(s, f"[{S(req_id)}] {req_title}", f"옵션 {sel} · {big}")

        meta = grp[grp["슬라이드번호"] == "META"]
        if not meta.empty:
            m = meta.iloc[0]
            bullets = []
            for raw in (m.get("왜_이_옵션"), m.get("적합_시그널"), m.get("리스크"), m.get("완화책")):
                line = first_meaningful_line(raw)
                if line:
                    bullets.append("• " + line)
            if bullets:
                add_textbox(s, 0.9, 2.2, 10.6, 0.5, f"옵션 {sel} · {big}", size=22, bold=True)
                bullets_from_paragraphs(s, 0.9, 3.0, 10.6, 3.8, bullets, size=body_size)

        detail = grp[grp["슬라이드번호"].apply(lambda v: S(v).isdigit())].copy()
        if not detail.empty:
            detail["슬라이드번호"] = detail["슬라이드번호"].astype(int)
            detail = detail.sort_values(by=["슬라이드번호"])
            for _, r in detail.iterrows():
                ss = prs.slides.add_slide(blank)
                add_title_subtitle(ss, S(r.get("제목")), f"옵션 {sel} · {big}")
                body_lines = lines_for_display(r.get("본문초안"))
                if body_lines:
                    bullets_from_paragraphs(ss, 0.9, 2.2, 10.6, 4.8, body_lines, size=body_size)
                urls = parse_url_list(r.get("URL"))
                if urls:
                    add_textbox(ss, 0.9, 7.3, 10.6, 0.5, "참고 URL", size=14, bold=True)
                    add_textbox(ss, 0.9, 7.8, 10.6, 0.8, "\n".join(urls), size=12)

    closing = prs.slides.add_slide(blank)
    add_title_subtitle(closing, "다음 단계", "")
    bullets_from_paragraphs(closing, 0.9, 2.2, 10.6, 3.0, [
        "옵션 선택 워크숍",
        "데이터/사전조건 점검",
        "파일럿 범위 합의 및 킥오프"
    ], size=body_size)

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()

# ==== 키 유틸 ====
import time, traceback, requests

import os, time, requests, traceback
import streamlit as st

# app.py 상단 (imports 아래쪽 어느 위치든 OK)
from dotenv import load_dotenv
load_dotenv(override=True)

import os, time, requests
from openai import OpenAI

def _mask(v: str | None, n: int = 8) -> str:
    return (v[:n] + "…") if v else ""

# 키 가져오기(환경변수/Secrets 모두 허용)
def _get_secret(key: str, default: str = "") -> str:
    val = os.getenv(key, default)
    try:
        import streamlit as st
        if not val:
            val = st.secrets.get(key, default)
    except Exception:
        pass
    return val

# 뱃지 HTML
def key_badge(key_name: str) -> str:
    val = _get_secret(key_name, "")
    state = "OK" if val else "NOT SET"
    color = "#10b981" if val else "#ef4444"
    tip = _mask(val) if val else "—"
    return f"""
    <span style="display:inline-block;padding:4px 8px;border-radius:999px;background:{color}22;
                 border:1px solid {color}44;color:{color};font-size:12px">
        {key_name}: <b>{state}</b> <span style="opacity:.7">({tip})</span>
    </span>
    """

# OpenAI 핑
def ping_openai(model: str = "gpt-4o-mini") -> dict:
    key = _get_secret("OPENAI_API_KEY")
    if not key:
        raise RuntimeError("OPENAI_API_KEY가 없습니다.")
    cli = OpenAI(api_key=key)
    t0 = time.time()
    resp = cli.chat.completions.create(
        model=model,
        messages=[{"role":"user","content":"ping"}],
        temperature=1,
    )
    lat = time.time() - t0
    txt = resp.choices[0].message.content.strip()
    return {"latency_sec": lat, "text": txt}

# Perplexity 핑
def _get_pplx_key() -> str:
    return (_get_secret("PERPLEXITY_API_KEY")
            or _get_secret("PPLX_API_KEY")
            or _get_secret("PEPLEXITY_API_KEY"))

def ping_perplexity(model: str = "sonar-pro") -> dict:
    key = _get_pplx_key()
    if not key:
        raise RuntimeError("Perplexity API 키가 없습니다. (PERPLEXITY_API_KEY / PPLX_API_KEY / PEPLEXITY_API_KEY)")
    url = "https://api.perplexity.ai/chat/completions"
    t0 = time.time()
    r = requests.post(
        url,
        headers={"Authorization": f"Bearer {key}", "Content-Type": "application/json"},
        json={
            "model": model,
            "messages": [{"role":"user","content":"ping"}],
            "temperature": 1
        },
        timeout=30
    )
    if r.status_code >= 400:
        raise RuntimeError(f"Perplexity {r.status_code}: {r.text}")
    lat = time.time() - t0
    data = r.json()
    txt = data["choices"][0]["message"]["content"].strip()
    return {"latency_sec": lat, "text": txt}


def build_excel(selected_df: pd.DataFrame) -> bytes:
    out = io.BytesIO()
    with pd.ExcelWriter(out, engine="xlsxwriter") as writer:
        selected_df.to_excel(writer, index=False, sheet_name="SelectedOptions")
        wb = writer.book
        ws = writer.sheets["SelectedOptions"]
        fmt = wb.add_format({"font_name": "Malgun Gothic", "font_size": 10})
        ws.set_column(0, selected_df.shape[1]-1, 24, fmt)
    return out.getvalue()

# =====================================================================================
# 인라인 미리보기 렌더러 (기존)
# =====================================================================================
def render_inline_preview(req_id: str, sub_df: pd.DataFrame, selected_opt: str):
    sub_df = sub_df.copy()
    req_title = S(sub_df["요청 제목"].iloc[0] if "요청 제목" in sub_df.columns and not sub_df.empty else req_id)
    opt_df = sub_df[sub_df["옵션번호"] == selected_opt]
    big = S(opt_df["옵션대제목"].iloc[0]) if not opt_df.empty and "옵션대제목" in opt_df.columns else ""

    st.markdown(f"**{req_title}**")
    st.caption(f"선택: 옵션 {selected_opt} — {big}")

    ov = sub_df[sub_df["슬라이드번호"] == "OVERVIEW"]
    if not ov.empty:
        ov_tab = try_extract_overview_table_from_row(ov.iloc[0])
        if ov_tab and ov_tab.get("columns") and ov_tab.get("rows"):
            with st.expander("옵션 비교(OVERVIEW)", expanded=False):
                st.table(pd.DataFrame(ov_tab["rows"], columns=ov_tab["columns"]))

    meta = opt_df[opt_df["슬라이드번호"] == "META"]
    if not meta.empty:
        m = meta.iloc[0]
        col1, col2 = st.columns(2)
        with col1:
            line = first_meaningful_line(m.get("왜_이_옵션"))
            if line:
                st.markdown("**왜 이 옵션인가**")
                st.write(line)
            line = first_meaningful_line(m.get("적합_시그널"))
            if line:
                st.markdown("**적합 시그널**")
                st.write(line)
        with col2:
            line = first_meaningful_line(m.get("리스크"))
            if line:
                st.markdown("**리스크**")
                st.write(line)
            line = first_meaningful_line(m.get("완화책"))
            if line:
                st.markdown("**완화책**")
                st.write(line)

        tl = parse_timeline(m.get("타임라인"))
        if tl:
            with st.expander("타임라인(주)", expanded=False):
                st.table(pd.DataFrame(tl))

    detail = opt_df[opt_df["슬라이드번호"].apply(lambda v: S(v).isdigit())].copy()
    if not detail.empty:
        detail["슬라이드번호"] = detail["슬라이드번호"].astype(int)
        detail = detail.sort_values("슬라이드번호")
        top_n = detail.head(2)
        for _, r in top_n.iterrows():
            st.markdown(f"- **{S(r.get('제목'))}**")
            fl = first_meaningful_line(r.get("본문초안"))
            if fl:
                st.write("  " + fl)

        with st.expander("상세 슬라이드 전체 보기", expanded=False):
            for _, r in detail.iterrows():
                st.markdown(f"**{S(r.get('제목'))}**")
                if S(r.get("부제목")):
                    st.caption(S(r.get("부제목")))
                body_lines = lines_for_display(r.get("본문초안"))
                for ln in body_lines:
                    st.write("- " + ln)
                urls = parse_url_list(r.get("URL"))
                if urls:
                    st.caption("참고 URL")
                    for u in urls:
                        st.write(u)
                st.divider()

# =====================================================================================
# Streamlit UI: 탭 구성 (1) Excel 기반 (2) RFP 업로드 (3) 히스토리
# =====================================================================================
st.set_page_config(page_title="Proposal Builder", layout="wide")
init_db()
login_form()
logout_button()

init_db()

st.title("제안 생성 · 미리보기 · 내보내기 (자동/엑셀)")
tab1, tab2, tab3 = st.tabs(["📥 Excel 업로드", "⚡ 자동 생성(LLM·RFP/방향성)", "🗂️ 고객 히스토리"])

# -------------------------- TAB 1: 기존 Excel 흐름 --------------------------
with tab1:
    if not is_authed():
        st.warning("이 메뉴는 로그인 후 이용할 수 있습니다.")
        st.stop()
    with st.sidebar:
        st.subheader("내보내기 설정")
        pdf_body_size = st.slider("PDF 본문 글자 크기", 9, 14, 11)
        ppt_body_size = st.slider("PPT 본문 글자 크기", 12, 20, 16)
        st.markdown("---")
        st.caption(f"OS: {platform.system()} | PDF Font: {PDF_KO_FONT or '기본'} | PPT Font: {_ppt_ko_font_name()}")

    st.markdown("#### 1) 데이터 업로드")
    uploaded = st.file_uploader("`slim_master_slide` CSV/Excel 업로드", type=["csv", "xlsx"])
    if uploaded is None:
        st.info("CSV/Excel를 업로드하세요. (필수 컬럼 예시: 요청 ID, 요청 제목, 옵션번호, 슬라이드번호, 제목, 부제목, 본문초안, 왜_이_옵션, 적합_시그널, 리스크, 완화책, 타임라인, URL, (선택) 옵션대제목)")
    else:
        if uploaded.name.lower().endswith(".csv"):
            df = pd.read_csv(uploaded, dtype=str).fillna("")
        else:
            df = pd.read_excel(uploaded, dtype=str).fillna("")

        required_cols = ["요청 ID","요청 제목","옵션번호","슬라이드번호","제목","부제목","본문초안","왜_이_옵션","적합_시그널","리스크","완화책","타임라인","URL"]
        for c in required_cols:
            if c not in df.columns:
                df[c] = ""

        df = compute_option_big_titles(df)
        df["슬라이드번호"] = df["슬라이드번호"].astype(str)
        df["옵션번호"] = df["옵션번호"].astype(str)

        st.markdown("#### 2) 고객 정보")
        col_a, col_b, col_c = st.columns(3)
        with col_a:
            client_name = st.text_input("고객사", value="")
        with col_b:
            author = st.text_input("작성팀", value="")
        with col_c:
            today_str = datetime.now().strftime("%Y-%m-%d")
            date_str = st.text_input("작성일", value=today_str)
        client_info = {"고객사": client_name, "작성팀": author, "작성일": date_str}

        st.markdown("#### 3) 요청별 옵션 선택 (아래에 즉시 미리보기)")
        req_ids = [x for x in df["요청 ID"].unique().tolist() if x not in ("COVER","CLOSING")]
        sel_map: Dict[str, str] = {}

        for rid in req_ids:
            sub = df[df["요청 ID"] == rid]
            req_title = S(sub["요청 제목"].iloc[0] if not sub.empty else rid)
            opts = sorted({o for o in sub["옵션번호"].unique().tolist() if S(o).isdigit()},
                          key=lambda x: int(x) if S(x).isdigit() else 999)
            if not opts:
                continue
            big_title_map = {}
            for o in opts:
                g = sub[sub["옵션번호"] == o]
                bt = S(g["옵션대제목"].iloc[0]) if not g.empty and "옵션대제목" in g.columns else ""
                big_title_map[o] = bt

            st.markdown(f"**[{rid}] {req_title}**")
            sel = st.radio(
                "옵션을 선택하세요",
                options=opts,
                horizontal=True,
                index=0,
                key=f"sel_{rid}",
                format_func=lambda o: f"{o} — {big_title_map.get(o, '')}" if big_title_map.get(o, "") else f"{o}"
            )
            sel_map[rid] = sel
            with st.container():
                render_inline_preview(rid, sub, sel)
            st.divider()

        frames = []
        cover_rows = df[df["요청 ID"] == "COVER"]
        if not cover_rows.empty:
            frames.append(cover_rows)
        for rid, sel in sel_map.items():
            sub = df[df["요청 ID"] == rid]
            overview = sub[sub["슬라이드번호"] == "OVERVIEW"]
            if not overview.empty:
                frames.append(overview)
            part = sub[sub["옵션번호"] == sel]
            frames.append(part)
        closing_rows = df[df["요청 ID"] == "CLOSING"]
        if not closing_rows.empty:
            frames.append(closing_rows)
        selected_df = pd.concat(frames, ignore_index=True) if frames else pd.DataFrame(columns=df.columns)

        st.markdown("#### 4) 내보내기")
        col1, col2, col3 = st.columns(3)
        with col1:
            if st.button("📄 PDF 생성", use_container_width=True):
                try:
                    pdf_bytes = build_pdf(selected_df, client_info, body_size=pdf_body_size)
                    st.success("PDF 생성 완료")
                    st.download_button("PDF 다운로드", data=pdf_bytes, file_name="proposal_options.pdf",
                                       mime="application/pdf", use_container_width=True)
                except Exception as e:
                    st.error(f"PDF 생성 오류: {e}")

        with col2:
            if st.button("📊 Excel 생성", use_container_width=True):
                try:
                    xlsx_bytes = build_excel(selected_df)
                    st.success("Excel 생성 완료")
                    st.download_button("Excel 다운로드", data=xlsx_bytes, file_name="proposal_options.xlsx",
                                       mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                                       use_container_width=True)
                except Exception as e:
                    st.error(f"Excel 생성 오류: {e}")

        with col3:
            if st.button("🖼️ PPT 생성", use_container_width=True):
                try:
                    ppt_bytes = build_ppt(selected_df, client_info, body_size=ppt_body_size)
                    st.success("PPT 생성 완료")
                    st.download_button("PPT 다운로드", data=ppt_bytes, file_name="proposal_options.pptx",
                                       mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                       use_container_width=True)
                except Exception as e:
                    st.error(f"PPT 생성 오류: {e}")

        with st.expander("선택 데이터 미리보기(전체)", expanded=False):
            st.dataframe(selected_df, use_container_width=True)

# ---- 자동 DF 생성 캐시 (같은 RFP/방향성 재실행 빠르게) ----
@st.cache_data(show_spinner=False, ttl=3600)
def _cached_build_flows_from_user_inputs(rfp_path, client_name, user_direction, notes, model_main):
    return build_flows_from_user_inputs(
        rfp_path=rfp_path,
        client_name=client_name,
        user_direction=user_direction,
        notes=notes,
        model_main=model_main,
    )


# ---------------- TAB 2: RFP 업로드/방향성 ----------------
# -------------------------- TAB 2: RFP 업로드/방향성 --------------------------
with tab2:
    if not is_authed():
        st.warning("이 메뉴는 로그인 후 이용할 수 있습니다.")
        st.stop()
    st.subheader("RFP 업로드 & 고객 방향성 입력")

    # ---- 세션 키들 ----
    LAST_PROJECT_KEY = "last_project"       # 최근 등록 프로젝트 메타
    AUTO_PAYLOAD_KEY = "auto_df_payload"    # 자동 생성 DF + 메타

    # 유틸: 스키마 보강
    def _ensure_slim_schema(df: pd.DataFrame) -> pd.DataFrame:
        need = ["요청 ID","요청 제목","옵션번호","슬라이드번호","제목","부제목","본문초안",
                "왜_이_옵션","적합_시그널","리스크","완화책","타임라인","URL","옵션대제목"]
        out = df.copy()
        for c in need:
            if c not in out.columns:
                out[c] = ""
        # astype errors='ignore' 는 일부 버전에서 경고를 내므로 try 처리
        try:
            out["슬라이드번호"] = out["슬라이드번호"].astype(str)
            out["옵션번호"] = out["옵션번호"].astype(str)
        except Exception:
            pass
        return out

    # 입력칸 초기화 제어
    if "rfp_form_version" not in st.session_state:
        st.session_state.rfp_form_version = 0
    def _reset_inputs():
        st.session_state.rfp_form_version += 1

    # 1) 고객 선택/등록
    with st.expander("고객 선택/등록", expanded=True):
        mode = st.radio("고객 선택 방식", ["기존 고객 선택", "신규 고객 등록"], horizontal=True,
                        key=f"mode_{st.session_state.rfp_form_version}")

        if mode == "기존 고객 선택":
            names = ["-- 선택 --"] + fetch_client_names()
            client_name = st.selectbox("고객명", options=names, index=0,
                                       key=f"client_select_{st.session_state.rfp_form_version}")
            capability = headcount = past_projects = ""
        else:
            client_name = st.text_input("고객명*", placeholder="예: 고객사명",
                                        key=f"client_new_{st.session_state.rfp_form_version}")
            capability = st.text_area("고객 역량/특기", key=f"cap_{st.session_state.rfp_form_version}")
            headcount = st.text_area("인원 정보", key=f"head_{st.session_state.rfp_form_version}")
            past_projects = st.text_area("전에 진행한 프로젝트", key=f"past_{st.session_state.rfp_form_version}")

    # 2) 프로젝트/방향성/RFP
    with st.form(f"rfp_form_{st.session_state.rfp_form_version}"):
        col1, col2 = st.columns(2)
        with col1:
            project_title = st.text_input("프로젝트 제목*", placeholder="예: 2025 이커머스 고도화 제안",
                                          key=f"title_{st.session_state.rfp_form_version}")
            direction = st.text_area("고객 방향성/원하는 바*", height=160,
                                     placeholder="예: 전환율 개선, CRM 연동, 보안 준수, 자동 리포트…",
                                     key=f"dir_{st.session_state.rfp_form_version}")
        with col2:
            rfp_file = st.file_uploader("RFP 파일 업로드 (PDF/DOCX/TXT)", type=["pdf", "docx", "txt"],
                                        key=f"rfp_{st.session_state.rfp_form_version}")
            notes = st.text_area("추가 메모(선택)", key=f"notes_{st.session_state.rfp_form_version}")

        a, b = st.columns([1,1])
        submitted = a.form_submit_button("등록하기")
        b.form_submit_button("입력칸 초기화", on_click=_reset_inputs)

    # 3) 제출 처리 → 세션에 저장하고 즉시 리런
    if submitted:
        if not client_name or client_name == "-- 선택 --":
            st.error("고객명을 선택/입력해주세요.")
            st.stop()
        if not project_title or not direction:
            st.error("프로젝트 제목과 고객 방향성은 필수입니다.")
            st.stop()
        if not rfp_file:
            st.error("RFP 파일을 업로드해주세요.")
            st.stop()

        client_id = upsert_client(client_name, capability, headcount, past_projects)
        project_id = create_project(client_id, project_title, direction)

        ts = _ts()
        safe_name = _sanitize(rfp_file.name)
        proj_rfp_dir = RFP_DIR / str(project_id)
        stored_path = proj_rfp_dir / f"{ts}__{safe_name}"
        _ensure_parent(stored_path)
        with open(stored_path, "wb") as f:
            f.write(rfp_file.getbuffer())
        attach_rfp(project_id, rfp_file.name, stored_path)

        proj_out_dir = RESULT_DIR / str(project_id)
        _ensure_parent(proj_out_dir / "config.json")
        config = {
            "project_id": project_id,
            "client_name": client_name,
            "project_title": project_title,
            "direction": direction,
            "rfp_path": str(stored_path),
            "notes": notes or "",
            "created_at": ts,
            "out_dir": str(proj_out_dir)
        }
        (proj_out_dir / "config.json").write_text(json.dumps(config, ensure_ascii=False, indent=2), encoding="utf-8")

        # ➜ 세션에 저장하고 리런 (아래 ‘자동 DF’ 섹션이 항상 보이게)
        st.session_state[LAST_PROJECT_KEY] = config
        st.session_state[AUTO_PAYLOAD_KEY] = None  # 새 프로젝트이므로 초기화
        st.rerun()

    # 4) 최근 프로젝트가 있으면 “자동 DF 생성” 섹션 항상 노출
    last_proj = st.session_state.get(LAST_PROJECT_KEY)

    # >>> AUTO_BUSY / 로그 초기화
    if "AUTO_BUSY" not in st.session_state:
        st.session_state.AUTO_BUSY = False
    if "autolog" not in st.session_state:
        st.session_state.autolog = []

    if last_proj:
        st.success(f"등록 완료! 프로젝트 ID: {last_proj['project_id']}")
        st.markdown("### ⚡ 엑셀 업로드 없이 바로 옵션 선택 (LLM 파이프라인)")

        # ▶ 눈에 확 띄는 상태 영역 + 진행바 + 로그
        status_box = st.empty()
        prog_bar = st.progress(0, text=("실행 중…" if st.session_state.AUTO_BUSY else "대기 중…"))
        log_expander = st.expander("실시간 진행 로그 보기", expanded=True)
        log_area = log_expander.empty()

        def _log(msg: str):
            # 진행 로그를 매번 갱신 렌더
            st.session_state.autolog.append(msg)
            log_area.write("\n".join(f"- {m}" for m in st.session_state.autolog[-200:]))

        # 버튼은 바쁠 때 비활성화
        has_openai = bool(os.getenv("OPENAI_API_KEY") or st.secrets.get("OPENAI_API_KEY", ""))
        has_pplx   = bool(os.getenv("PERPLEXITY_API_KEY") or st.secrets.get("PERPLEXITY_API_KEY", "") or
                        os.getenv("PEPLEXITY_API_KEY") or os.getenv("PPLX_API_KEY"))
        btn_disabled = (not has_openai) or (not has_pplx) or st.session_state.AUTO_BUSY
        btn_label = "⏳ 생성 중…" if st.session_state.AUTO_BUSY else "🚀 자동 DF 생성 (고객 방향성 + RFP 기반)"

        action_col1, _ = st.columns([1, 3])
        with action_col1:
            if st.button(btn_label, use_container_width=True,
                        key=f"autodf_{last_proj['project_id']}", disabled=btn_disabled):

                # 필수 키 점검(명확한 에러)
                if not has_openai:
                    st.error("OpenAI API 키가 없습니다. .streamlit/secrets.toml 또는 환경변수에 설정하세요.")
                    st.stop()
                if not has_pplx:
                    st.error("Perplexity API 키가 없습니다. .streamlit/secrets.toml 또는 환경변수에 설정하세요.")
                    st.stop()

                # >>> 바쁨 상태 시작 (중복 실행 방지)
                st.session_state.AUTO_BUSY = True
                # 기존 로그/프로그레스 초기화
                st.session_state.autolog = []
                prog_bar.progress(0, text="초기화…")

                # 큰 상태 패널
                with status_box.status("LLM 준비 중…", expanded=True) as st_status:
                    try:
                        prog_bar.progress(5, text="세션/환경 초기화")
                        _log("세션/환경 초기화 완료")

                        st.write("1/4 RFP 텍스트 파싱…")
                        _log("RFP 텍스트 추출 시작")
                        prog_bar.progress(15, text="RFP 텍스트 추출…")
                        try:
                            if 'extract_text_from_file' in globals():
                                _ = extract_text_from_file(last_proj["rfp_path"])
                                _log("RFP 텍스트 추출 완료")
                            else:
                                _log("RFP 텍스트 추출 스킵(함수 미로딩)")
                        except Exception as e:
                            _log(f"RFP 텍스트 추출 경고: {e}")


                        st.write("2/4 LLM 분석(요청/질문/업데이트 플랜)…")
                        prog_bar.progress(35, text="LLM 분석 시작…")
                        _log("OpenAI/Perplexity 호출 시작")

                        # 핵심: 내부에서 LLM 호출 (시간 걸림) — 진행 중 UI가 유지됨
                        with st.spinner("🚀 LLM 분석 및 슬라이드 흐름 생성 중…"):
                            # UI 로그 함수를 가급적 전달하되, 시그니처에 없으면 어댑터가 제거
                            auto_df = _call_build_flows_adapted(
                                build_flows_from_user_inputs,
                                rfp_path=last_proj["rfp_path"],
                                client_name=last_proj["client_name"],
                                user_direction=last_proj["direction"],
                                notes=last_proj.get("notes", ""),
                                model_main="gpt-5",
                                model_deck="gpt-5",
                                out_dir=last_proj["out_dir"],
                                logf=_log,     # ← 있으면 전달, 없으면 자동 필터
                            )

                        prog_bar.progress(70, text="스키마 정리…")

                        st.write("3/4 스키마/옵션 제목 정리…")
                        auto_df = _ensure_slim_schema(auto_df)
                        auto_df = compute_option_big_titles(auto_df)
                        _log("스키마/옵션대제목 정리 완료")

                        st.write("4/4 결과 저장…")
                        out_dir = Path(last_proj["out_dir"])
                        out_dir.mkdir(parents=True, exist_ok=True)
                        auto_xlsx = out_dir / f"auto_df_{last_proj['project_id']}.xlsx"
                        with pd.ExcelWriter(auto_xlsx, engine="xlsxwriter") as w:
                            auto_df.to_excel(w, index=False, sheet_name="auto_df")
                        _log(f"파일 저장: {auto_xlsx}")

                        # 세션에 결과 탑재
                        st.session_state[AUTO_PAYLOAD_KEY] = {
                            "df": auto_df,
                            "meta": {
                                "project_id": last_proj["project_id"],
                                "client_name": last_proj["client_name"],
                                # created_at 키가 없을 수도 있어 호환 처리
                                "created_at": last_proj.get("created_at") or last_proj.get("ts", ""),
                            }
                        }
                        prog_bar.progress(100, text="완료")
                        st_status.update(label="✅ 자동 DF 생성 완료 — 아래에서 옵션을 선택하세요.", state="complete")
                        st.toast("자동 DF 생성 완료!", icon="✅")
                        try:
                            st.rerun()
                        except Exception:
                            pass  # Streamlit 내부 RerunData 예외는 무시

                    except Exception as e:
                        prog_bar.progress(100, text="오류")
                        _log(f"오류 발생: {e}")
                        st_status.update(label="❌ 실패", state="error")
                        st.error(f"자동 DF 생성 실패: {e}")
                        # 필요 시 디버그용
                        # st.code(traceback.format_exc())
                        # >>> 바쁨 상태 해제
                        st.session_state.AUTO_BUSY = False

        # 5) 자동 DF가 있으면 즉시 Excel 업로드와 동일한 플로우로 진행
        payload = st.session_state.get(AUTO_PAYLOAD_KEY)
        if payload and "df" in payload:
            auto_df: pd.DataFrame = payload["df"]
            meta = payload.get("meta", {})
            st.markdown("---")
            st.markdown("### ✅ 자동 생성 DF — Excel 업로드와 동일한 흐름으로 진행")

            # Tab1과 동일한 공용 컴포넌트(옵션 선택 · 미리보기 · 내보내기) 한 번만 렌더
            render_excel_like_flow(
                auto_df,
                default_client = meta.get("client_name", ""),
                default_author = "",
                default_date   = datetime.now().strftime("%Y-%m-%d"),
                pdf_body_size  = st.session_state.get("pdf_body_size", 11),   # ← 사이드바 값 재사용
                ppt_body_size  = st.session_state.get("ppt_body_size", 16),   # ← 사이드바 값 재사용
                key_prefix     = f"auto_{meta.get('project_id','X')}"         # ← 프로젝트별 유니크 키
            )

            with st.expander("자동 생성 DF(전체) 미리보기", expanded=False):
                st.dataframe(auto_df, use_container_width=True)

            # 숨기기 버튼도 한 번만 · 유니크 키
            if st.button(
                "🧹 이 자동 DF 숨기기/초기화",
                help="세션에서 제거합니다. 파일은 유지",
                key=f"clear_auto_df_{meta.get('project_id','X')}"
            ):
                st.session_state[AUTO_PAYLOAD_KEY] = None
                st.toast("자동 DF를 숨겼습니다.", icon="🧹")
                st.rerun()


# -------------------------- TAB 3: 고객/프로젝트 히스토리 --------------------------
# -------------------------- TAB 3: 고객/프로젝트 히스토리 --------------------------
with tab3:
    # 로그인 보호 (secrets.auth.enabled=true 인 경우만 동작)
    if not is_authed():
        st.warning("이 메뉴는 로그인 후 이용할 수 있습니다.")
        st.stop()

    st.subheader("고객/프로젝트 히스토리")

    names = ["-- 선택 --"] + fetch_client_names()
    selected_name = st.selectbox("고객 선택", options=names, index=0)

    if selected_name and selected_name != "-- 선택 --":
        client_id = fetch_client_id_by_name(selected_name)
        cap, head, past = fetch_client_info(client_id)

        st.markdown("### 고객 정보")
        st.write("**역량**:", cap or "-")
        st.write("**인원 정보**:", head or "-")
        st.write("**이전 프로젝트**:", past or "-")

        st.divider()
        st.markdown("### 프로젝트 목록")

        # 고객 전체 삭제
        if st.button("⚠️ 이 고객 전체 삭제", key=f"del_client_{client_id}", help="모든 프로젝트/파일이 삭제됩니다. 복구 불가"):
            delete_client_and_all(client_id)
            st.toast("고객 및 모든 프로젝트 삭제 완료", icon="🗑")
            st.rerun()

        projects = list_projects(client_id)
        if not projects:
            st.info("등록된 프로젝트가 없습니다.")
        else:
            for pid, title, direction, created_at, status in projects:
                row = st.container()
                left, right = row.columns([0.88, 0.12])

                with left:
                    with st.expander(f"[#{pid}] {title} — {status} ({created_at})", expanded=False):
                        # 기본 메타
                        st.markdown("**방향성**")
                        st.write(direction or "-")
                        st.markdown("---")

                        # ✅ 전체 결과 내역(다운로드/미리보기) 표시
                        _show_project_full_result(pid, key_prefix=f"proj_{pid}")

                with right:
                    st.write("")  # vertical spacing
                    st.write("")
                    if st.button("🗑 삭제", key=f"del_proj_{pid}", help="이 프로젝트와 관련 파일을 모두 삭제합니다."):
                        delete_project(pid)
                        st.toast(f"프로젝트 #{pid} 삭제 완료", icon="🗑")
                        st.rerun()
